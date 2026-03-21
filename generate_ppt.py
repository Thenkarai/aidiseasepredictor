from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def set_slide_background(slide, color_rgb):
    """Sets the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Adds an image to the slide if the file exists."""
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)

def create_presentation():
    # Paths to relevant images
    MOCKUP_PATH = r"C:\Users\Vichu\.gemini\antigravity\brain\0652770b-dac0-4c6b-82f2-8b073b959e3e\phasal_ui_mockup_1773801180629.png"
    CROP_IMG = r"c:\Users\Vichu\OneDrive\Documents\Plant-Disease-Recognition-System-main\aidiseasepredictor\static\images\crop_background.jpg"
    
    prs = Presentation()
    
    # Define Theme Colors
    BG_COLOR = RGBColor(10, 10, 25)      # Near Black / Deep Space
    ACCENT_COLOR = RGBColor(39, 174, 96) # Forest Green / Emerald
    TEXT_COLOR = RGBColor(240, 240, 240) # Bright White
    
    def add_styled_slide(title_text, bullet_points=None, image_path=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, BG_COLOR)
        
        # Title Styling
        title = slide.shapes.title
        title.text = title_text
        for p in title.text_frame.paragraphs:
            p.font.name = "Impact"
            p.font.size = Pt(40)
            p.font.color.rgb = ACCENT_COLOR
        
        # Content Styling
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        if bullet_points:
            # Adjust body width if image is present
            if image_path:
                body_shape.width = Inches(5.5)
            
            for point in bullet_points:
                p = tf.add_paragraph()
                p.text = "  • " + point
                p.font.name = "Segoe UI"
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(10)

        if image_path:
            add_image_to_slide(slide, image_path, Inches(6), Inches(1.5), width=Inches(3.5))

        # Bottom accent bar
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(7.1), Inches(9), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.visible = False

    # Slide 1: Professional Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "PHASAL"
    title_p = title.text_frame.paragraphs[0]
    title_p.font.name = "Impact"
    title_p.font.size = Pt(88)
    title_p.font.color.rgb = ACCENT_COLOR
    
    subtitle.text = "AI Plant Disease Recognition System\n\nPresented by:\nThenkarai Maharajan M\nMCA Student"
    for p in subtitle.text_frame.paragraphs:
        p.font.name = "Segoe UI Semibold"
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    # Slide 2: Project Overview (with Image)
    add_styled_slide("Introduction", [
        "Revolutionizing Agriculture via AI.",
        "Real-time Identification of plant pathologies.",
        "Bridging the Gap between technology and farmers.",
        "Comprehensive health reports in English and Tamil."
    ], image_path=CROP_IMG)

    # Slide 3: Problem Statement
    add_styled_slide("Problem Statement", [
        "Global crop yield loss (20-40%) due to diseases.",
        "Limited access to agriculture experts.",
        "Slow manual diagnostics leading to crop failure.",
        "Need for localized linguistic support (Tamil)."
    ])

    # Slide 4: System Objectives
    add_styled_slide("Project Objectives", [
        "Develop a lightweight, high-accuracy AI model.",
        "Provide Instant analysis via Web/Camera.",
        "Quantify Severity for better decision making.",
        "Enable Tamil support for regional accessibility."
    ])

    # Slide 5: Tech Stack
    add_styled_slide("Technology Stack", [
        "Backend: Python / Flask Web Server.",
        "AI: Google Gemini-2.5-Flash (LLM Vision).",
        "Frontend: Modern HTML5 / CSS3 / JavaScript.",
        "Processing: PIL (Pillow) & Base64 Logic.",
        "UI Theme: Professional Dark Mode & Glassmorphism."
    ])

    # Slide 6: User Interface (The "Output Image" Slide)
    add_styled_slide("User Interface & Analysis", [
        "Clean, intuitive dashboard for farmers.",
        "Real-time confidence circular gauges.",
        "Interactive disease severity indicators.",
        "Localized result cards for quick reading."
    ], image_path=MOCKUP_PATH)

    # Slide 7: System Workflow
    add_styled_slide("System Workflow", [
        "1. Image Capture: User uploads or uses camera.",
        "2. Analysis: Gemini AI processes visual features.",
        "3. Quantification: Severity & Confidence calculated.",
        "4. Reporting: Results displayed with Cause & Cure.",
        "5. Localization: Dynamic translation to Tamil."
    ])

    # Slide 8: Key Features
    add_styled_slide("Key Features", [
        "Multi-modal AI: Analyzes images with high precision.",
        "Tamil Translation: Zero barrier for local farmers.",
        "Severity Analysis: Precise leaf area impact data.",
        "Cross-platform: Works on mobile and desktop."
    ])

    # Slide 9: Future Scope
    add_styled_slide("Conclusion & Future Scope", [
        "Integration with IoT Soil Sensors.",
        "Development of Dedicated Mobile Apps.",
        "Offline Database for remote rural areas.",
        "Blockchain for secure crop data tracking."
    ])

    # Slide 10: Final Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, BG_COLOR)
    title = slide.shapes.title
    title.text = "Thank You!"
    title_p = title.text_frame.paragraphs[0]
    title_p.font.name = "Impact"
    title_p.font.size = Pt(80)
    title_p.font.color.rgb = ACCENT_COLOR
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Questions & Feedback\n\nThenkarai Maharajan M\nMCA Student"
    for p in subtitle.text_frame.paragraphs:
        p.font.name = "Segoe UI"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR

    prs.save("PHASAL_Final_Presentation.pptx")
    print("Professional MCA PPT created successfully: PHASAL_Final_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
